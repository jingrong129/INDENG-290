from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = PROJECT_ROOT / "outputs" / "oil_forecast_evaluation_presentation.pptx"

# ── Palette ────────────────────────────────────────────────────────────────
NAVY = RGBColor(20, 34, 54)          # primary text
BLUE = RGBColor(37, 99, 145)         # secondary accent
TEAL = RGBColor(45, 139, 132)        # primary accent / positive
ORANGE = RGBColor(206, 116, 43)      # risk / bias callouts
RED = RGBColor(170, 65, 62)          # strong warning
LIGHT_BG = RGBColor(248, 249, 250)   # off-white slide background
MID_GRAY = RGBColor(108, 118, 130)   # metadata / subtitles
LIGHT_GRAY = RGBColor(218, 226, 232) # dividers / borders
WHITE = RGBColor(255, 255, 255)
TEAL_BG = RGBColor(230, 242, 241)    # very light teal fill
BLUE_BG = RGBColor(232, 242, 252)    # very light blue fill
ORANGE_BG = RGBColor(252, 242, 234)  # very light orange fill


# ── Low-level text helpers ─────────────────────────────────────────────────

def _para_run(tf, text, font_size, color, bold=False, align=PP_ALIGN.LEFT,
              first=False, space_before=0, line_spacing=None):
    """Add a paragraph+run to a text frame. Uses first paragraph when first=True."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def _tf_setup(tf, word_wrap=True, vertical_anchor=MSO_ANCHOR.MIDDLE,
              margin_h=Inches(0.06), margin_v=Inches(0.05)):
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = margin_h
    tf.margin_right = margin_h
    tf.margin_top = margin_v
    tf.margin_bottom = margin_v


def add_textbox(slide, left, top, width, height, text, font_size=22,
                color=NAVY, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.65), Inches(0.28), Inches(11.8), Inches(0.62),
                title, 28, NAVY, True)
    if subtitle:
        add_textbox(slide, Inches(0.65), Inches(0.88), Inches(11.5), Inches(0.30),
                    subtitle, 13, MID_GRAY)
    # Short teal accent bar under title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.65), Inches(1.24), Inches(1.6), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_footer(slide, number):
    add_textbox(slide, Inches(11.85), Inches(7.12), Inches(0.85), Inches(0.22),
                str(number), 10, LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def add_divider(slide, left, top, width, color=LIGHT_GRAY):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_bullets(slide, items, left, top, width, height, font_size=18,
                color=NAVY, gap_before=6, bullet_char="•"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if idx:
            p.space_before = Pt(gap_before)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_table(slide, data, left, top, width, height, col_widths=None,
              font_size=14, header_fill=NAVY):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = cw
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.07)
            cell.margin_bottom = Inches(0.07)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            paragraph.line_spacing = 1.1
            run = paragraph.runs[0]
            run.font.name = "Aptos"
            run.font.size = Pt(font_size)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                run.font.bold = True
                run.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT_BG
                run.font.color.rgb = NAVY
    return table_shape


def add_picture_fit(slide, image_path, left, top, width, height):
    from PIL import Image
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    target_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > target_ratio:
        pic_w = width
        pic_h = width / img_ratio
    else:
        pic_h = height
        pic_w = height * img_ratio
    return slide.shapes.add_picture(
        str(image_path),
        left + (width - pic_w) / 2,
        top + (height - pic_h) / 2,
        pic_w, pic_h,
    )


def set_notes(slide, notes):
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    frame.text = notes


def add_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_title(slide, title, subtitle)
    add_footer(slide, len(prs.slides))
    return slide


# ── Composite element helpers ──────────────────────────────────────────────

def _shape_text(slide, shape_type, x, y, w, h, fill_color, line_color,
                paragraphs, line_width_pt=1.2, vertical=MSO_ANCHOR.MIDDLE,
                margin_h=Inches(0.14), margin_v=Inches(0.1)):
    """Add a filled shape with one or more paragraphs of text.

    paragraphs: list of (text, font_size, color, bold, align, space_before, line_spacing)
    """
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    tf = shape.text_frame
    _tf_setup(tf, vertical_anchor=vertical, margin_h=margin_h, margin_v=margin_v)
    for i, para in enumerate(paragraphs):
        text, fs, col, bold, align = para[:5]
        space_before = para[5] if len(para) > 5 else 0
        ls = para[6] if len(para) > 6 else 1.15
        _para_run(tf, text, fs, col, bold, align, first=(i == 0),
                  space_before=space_before, line_spacing=ls)
    return shape


def add_flow_node(slide, x, y, w, h, heading, body_lines=None):
    paras = [(heading, 15, TEAL, True, PP_ALIGN.CENTER)]
    if body_lines:
        paras.append((body_lines, 12.5, NAVY, False, PP_ALIGN.CENTER, 5, 1.1))
    _shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                WHITE, LIGHT_GRAY, paras, line_width_pt=1.0,
                vertical=MSO_ANCHOR.TOP, margin_v=Inches(0.1))


def add_number_callout(slide, x, y, w, label, number, detail,
                       number_color=TEAL, label_size=12, number_size=54, detail_size=14):
    """Stacked large-number callout: label / big number / detail."""
    add_textbox(slide, x, y, w, Inches(0.30), label.upper(), label_size, MID_GRAY,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.28), w, Inches(1.05), number,
                number_size, number_color, True, PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(1.32), w, Inches(0.35), detail,
                detail_size, MID_GRAY, align=PP_ALIGN.CENTER)


# ── Main deck builder ──────────────────────────────────────────────────────

def build_deck():
    benchmark = pd.read_csv(PROJECT_ROOT / "data/processed/benchmark_metrics.csv")
    rolling = pd.read_csv(PROJECT_ROOT / "data/processed/rolling_model_metrics.csv")

    def metric(model, horizon, column, frame=benchmark):
        row = frame.loc[
            (frame["model"] == model) & (frame["horizon_years"] == horizon)
        ].iloc[0]
        return float(row[column])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 1: TITLE ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    # Top teal bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL
    top_bar.line.fill.background()
    # Title
    add_textbox(slide, Inches(0.78), Inches(1.75), Inches(10.5), Inches(1.05),
                "Long-Term Oil Price Forecast Evaluation",
                38, WHITE, True, PP_ALIGN.LEFT)
    # Teal accent rule under title
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(2.82), Inches(2.8), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()
    # Subtitle
    add_textbox(slide, Inches(0.82), Inches(2.98), Inches(9.8), Inches(0.56),
                "Are institutional Brent forecasts useful for upstream investment planning?",
                18, RGBColor(190, 212, 228), False, PP_ALIGN.LEFT)
    # Decorative bars (subtle bar-chart motif, bottom-right)
    for bx, bh, bc in [(9.05, 1.8, BLUE), (9.6, 2.7, TEAL), (10.15, 1.3, ORANGE), (10.7, 2.2, RED)]:
        b = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bx), Inches(5.78 - bh), Inches(0.36), Inches(bh))
        b.fill.solid()
        b.fill.fore_color.rgb = bc
        b.line.fill.background()
    # Bottom metadata
    add_textbox(slide, Inches(0.82), Inches(6.62), Inches(8.0), Inches(0.32),
                "Berkeley Energy Project  |  Annual Brent Spot Price in real 2025 USD/bbl",
                11.5, RGBColor(90, 115, 140), align=PP_ALIGN.LEFT)
    set_notes(slide, (
        "Good afternoon. I'm [name], and today I'll be presenting our evaluation of EIA long-term "
        "Brent crude oil price forecasts. The central question is whether these institutional "
        "forecasts are actually useful for upstream investment planning — or whether their errors "
        "are large enough to undermine the capital allocation decisions that depend on them. "
        "Over the next twenty minutes, I'll walk you through our data, our methodology, and two "
        "key findings with direct implications for how energy firms should be using these price decks."
    ))

    # ── 2: DECISION CONTEXT ───────────────────────────────────────────────
    slide = add_slide(prs, "Decision Context",
                      "Why forecast accuracy matters for upstream oil and gas firms")
    # Intro line
    add_textbox(slide, Inches(0.72), Inches(1.42), Inches(11.9), Inches(0.40),
                "Upstream capital projects are large, long-lived, and partially irreversible. "
                "Oil price assumptions flow directly into DCF valuations.",
                17, NAVY, align=PP_ALIGN.LEFT)
    # Flow nodes
    flow_items = [
        ("Oil Price\nForecast", None),
        ("Revenue\nProjection", None),
        ("DCF / NPV\n/ IRR", None),
        ("Investment\nDecision", None),
    ]
    node_w = Inches(2.1)
    node_h = Inches(1.05)
    gap = Inches(0.52)
    total_w = len(flow_items) * node_w + (len(flow_items) - 1) * gap
    start_x = (Inches(13.333) - total_w) / 2
    flow_y = Inches(2.22)
    for i, (text, _) in enumerate(flow_items):
        nx = start_x + i * (node_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx, flow_y, node_w, node_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL_BG
        shape.line.color.rgb = TEAL
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        _tf_setup(tf)
        _para_run(tf, text, 16, NAVY, True, PP_ALIGN.CENTER, first=True, line_spacing=1.1)
        if i < len(flow_items) - 1:
            add_textbox(slide, nx + node_w + Inches(0.09), flow_y + Inches(0.34),
                        Inches(0.34), Inches(0.38), "→", 22, TEAL, True, PP_ALIGN.CENTER)
    # Risk callout
    add_divider(slide, Inches(0.72), Inches(3.55), Inches(11.9), ORANGE)
    add_textbox(slide, Inches(0.72), Inches(3.68), Inches(11.9), Inches(0.50),
                "⚠   Risk: Systematic overprediction encourages overinvestment and future asset impairments.",
                19, ORANGE, True, PP_ALIGN.CENTER)
    # Supporting bullets
    add_bullets(slide,
                ["Large upfront capital commitments with multi-year payoff horizons",
                 "Price assumptions drive revenue projections, NPV, and IRR",
                 "Partially irreversible decisions — cost of over-approval is high"],
                Inches(1.5), Inches(4.38), Inches(10.3), Inches(1.8),
                18, NAVY, gap_before=7)
    set_notes(slide, (
        "Let me start with why this matters beyond forecasting statistics. "
        "Upstream oil and gas investments are large, long-lived, and largely irreversible. "
        "When a firm sanctions a deepwater field or commits to a major infrastructure project, "
        "it is committing capital that cannot easily be recovered if things go wrong.\n\n"
        "The chain of logic on this slide is the one used in every upstream investment decision. "
        "An oil price forecast feeds into a revenue projection. That projection goes into a "
        "discounted cash flow model to compute NPV and IRR. Based on those metrics, the "
        "investment committee approves or rejects.\n\n"
        "Here is the problem. If the price forecast feeding that DCF is systematically too high "
        "— if it is optimistically biased — then every downstream number is inflated. Revenue "
        "looks better than it should. NPV looks better. IRR clears the hurdle rate more easily. "
        "And the firm ends up approving projects it should not have, or approving them at too "
        "large a scale.\n\n"
        "That is the risk flagged at the bottom of this slide. Systematic overprediction creates "
        "a pipeline of overvalued projects that look attractive in the model but disappoint when "
        "realized prices arrive. This eventually shows up as write-downs and asset impairments.\n\n"
        "The quality of long-term price forecasts is therefore not an academic question. "
        "It has direct consequences for capital allocation — which is what makes this evaluation "
        "worth doing."
    ))

    # ── 3: RESEARCH QUESTIONS ─────────────────────────────────────────────
    slide = add_slide(prs, "Research Questions")
    questions = [
        ("1", TEAL,
         "Does EIA beat a naive benchmark?",
         "Does the institutional forecast outperform a simple random-walk carry-forward rule?"),
        ("2", TEAL,
         "Is there systematic forecast bias?",
         "Does EIA consistently overpredict or underpredict realized annual Brent prices?"),
        ("3", BLUE,
         "Can simple models improve on EIA?",
         "Do rolling AR models with or without futures data reduce forecast error further?"),
    ]
    row_gap = Inches(1.58)
    for i, (num, col, q, detail) in enumerate(questions):
        y = Inches(1.52) + i * row_gap
        # Numbered circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(0.72), y + Inches(0.08), Inches(0.62), Inches(0.62))
        circ.fill.solid()
        circ.fill.fore_color.rgb = col
        circ.line.fill.background()
        ctf = circ.text_frame
        _tf_setup(ctf, margin_h=Inches(0.0), margin_v=Inches(0.0))
        _para_run(ctf, num, 18, WHITE, True, PP_ALIGN.CENTER, first=True)
        # Question (large)
        add_textbox(slide, Inches(1.58), y, Inches(11.0), Inches(0.46),
                    q, 21, NAVY, True, PP_ALIGN.LEFT)
        # Detail (smaller gray)
        add_textbox(slide, Inches(1.58), y + Inches(0.46), Inches(11.0), Inches(0.38),
                    detail, 15, MID_GRAY, False, PP_ALIGN.LEFT)
        if i < 2:
            add_divider(slide, Inches(0.72), y + Inches(1.08), Inches(11.95))
    # Scope note at bottom
    add_textbox(slide, Inches(0.72), Inches(6.38), Inches(11.9), Inches(0.38),
                "Horizons evaluated: 3-year and 5-year ahead forecasts of annual average Brent Spot Price in real 2025 USD/bbl.",
                13.5, MID_GRAY, align=PP_ALIGN.CENTER)
    set_notes(slide, (
        "Given that context, our study is organized around three research questions.\n\n"
        "First: does EIA beat a naive benchmark? Before we can say a forecast is useful, we need "
        "to show it adds value beyond the simplest possible alternative. If EIA cannot beat a "
        "carry-forward rule, there is no empirical case for preferring it.\n\n"
        "Second: is there systematic bias? A forecast can have reasonable overall accuracy and "
        "still be consistently too high or too low. We test specifically for directional error "
        "that would affect investment decisions.\n\n"
        "Third: can simple models improve on EIA? We built rolling autoregressive models with and "
        "without futures price data, and tested whether they could reduce forecast error further.\n\n"
        "Throughout, we evaluate three-year and five-year ahead forecasts of annual average "
        "Brent Spot Price in real 2025 U.S. dollars per barrel."
    ))

    # ── 4: DATA AND SCOPE ─────────────────────────────────────────────────
    slide = add_slide(prs, "Data and Scope",
                      "Building a reproducible forecast vintage panel")
    sources = [
        ("EIA AEO Table 12",
         "Brent forecast vintages — primary evaluation target",
         "Clean annual AEO releases from 2013–2023; no 2024 AEO was published."),
        ("FRED Brent  (DCOILBRENTEU)",
         "Realized Brent price — outcome variable",
         "Daily spot prices aggregated to annual averages; CPI-deflated to real 2025 USD/bbl."),
        ("FRED CPI  (CPIAUCSL)",
         "Inflation deflator",
         "Converts nominal EIA vintage forecasts to real 2025 USD/bbl for consistent comparison."),
        ("Yahoo Finance BZ=F",
         "Front-month Brent futures — market signal",
         "Used as an additional regressor in the rolling AR + Futures extension."),
    ]
    row_h = Inches(1.14)
    for i, (src, var, note) in enumerate(sources):
        y = Inches(1.40) + i * row_h
        add_textbox(slide, Inches(0.72), y, Inches(3.55), Inches(0.38),
                    src, 15.5, TEAL, True, PP_ALIGN.LEFT)
        add_textbox(slide, Inches(4.35), y, Inches(8.3), Inches(0.38),
                    var, 15.5, NAVY, True, PP_ALIGN.LEFT)
        add_textbox(slide, Inches(0.72), y + Inches(0.38), Inches(12.0), Inches(0.36),
                    note, 13.5, MID_GRAY, False, PP_ALIGN.LEFT)
        if i < 3:
            add_divider(slide, Inches(0.72), y + Inches(0.82), Inches(11.9))
    # Footer note on World Bank
    add_textbox(slide, Inches(0.72), Inches(6.22), Inches(11.9), Inches(0.32),
                "Note: World Bank Crude oil (avg) used as supplementary check only — different contract from Brent spot.",
                13, MID_GRAY, align=PP_ALIGN.LEFT)
    set_notes(slide, (
        "Our data comes from four main sources.\n\n"
        "The primary forecast data is EIA's Annual Energy Outlook, Table 12, which reports "
        "long-term Brent price projections each year. We collected AEO releases from 2013 through "
        "2023. There was no AEO published in 2024, so that year is excluded. This gives us "
        "eleven distinct vintage forecasts to evaluate.\n\n"
        "For realized Brent prices, we use FRED's daily Brent spot price series — "
        "DCOILBRENTEU — aggregated to annual averages. This is our outcome variable: what "
        "actually happened.\n\n"
        "Because different AEO vintages report prices in different base years, we cannot directly "
        "compare across vintages without standardizing. We use the FRED CPI series — CPIAUCSL "
        "— to convert all values to real 2025 U.S. dollars per barrel. This is a critical step "
        "that is easy to get wrong.\n\n"
        "For the rolling model extension we will describe later, we use the Yahoo Finance BZ=F "
        "front-month Brent futures contract as a market-based signal.\n\n"
        "One note: World Bank crude prices are available but track a composite crude average "
        "rather than Brent specifically, so we treat them as supplementary rather than core data."
    ))

    # ── 5: EVALUATION PIPELINE ────────────────────────────────────────────
    slide = add_slide(prs, "Evaluation Pipeline",
                      "Standardizing forecasts and realized prices into one panel")
    steps = [
        ("EIA Vintages", "AEO Table 12\nforecast releases"),
        ("Realized Brent", "FRED daily spot\n→ annual average"),
        ("CPI Deflation", "Nominal →\nReal 2025 USD/bbl"),
        ("Panel", "3y & 5y horizon\noutcome pairs"),
        ("Benchmark", "Random-walk\ncomparison"),
    ]
    n = len(steps)
    node_w = Inches(1.88)
    node_h = Inches(1.45)
    arr_w = Inches(0.44)
    total = n * node_w + (n - 1) * arr_w
    sx = (Inches(13.333) - total) / 2
    ny = Inches(2.05)
    for i, (heading, body) in enumerate(steps):
        nx = sx + i * (node_w + arr_w)
        add_flow_node(slide, nx, ny, node_w, node_h, heading, body)
        if i < n - 1:
            add_textbox(slide, nx + node_w + Inches(0.04), ny + Inches(0.50),
                        Inches(0.38), Inches(0.45), "→", 22, TEAL, True, PP_ALIGN.CENTER)
    # Stat callouts below pipeline
    stats = [
        ("EIA Publication Vintages", "11", "evaluation windows"),
        ("Forecast-Outcome Pairs", "20", "3-year and 5-year"),
        ("Price Unit", "2025 $", "real USD / bbl"),
    ]
    for j, (lbl, val, sub) in enumerate(stats):
        cx = Inches(1.65) + j * Inches(3.42)
        add_textbox(slide, cx, Inches(3.85), Inches(3.1), Inches(0.30),
                    lbl, 12, MID_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx, Inches(4.12), Inches(3.1), Inches(0.78),
                    val, 40, TEAL, True, PP_ALIGN.CENTER)
        add_textbox(slide, cx, Inches(4.88), Inches(3.1), Inches(0.28),
                    sub, 12, MID_GRAY, align=PP_ALIGN.CENTER)
    set_notes(slide, (
        "Let me walk through how we turn those raw data sources into an evaluation panel.\n\n"
        "We start with EIA vintage data from AEO Table 12. Each vintage gives us a multi-year "
        "forecast path. We extract the price predicted for three and five years ahead from each "
        "release date.\n\n"
        "In parallel, we compute annual average Brent prices from FRED's daily spot series. "
        "These are the realized outcomes we compare against.\n\n"
        "The CPI deflation step converts everything to real 2025 dollars. This is essential "
        "because EIA's own base years shift across vintages — without this step, a dollar "
        "comparison across years would be meaningless.\n\n"
        "After matching each forecast to its realized outcome, we have our evaluation panel: "
        "forecast-outcome pairs at three-year and five-year horizons for each vintage year. "
        "That is twenty observation pairs in total.\n\n"
        "The final piece is the random-walk benchmark, constructed in parallel. For each vintage, "
        "the random walk forecast is simply the most recent completed annual Brent average "
        "at the time of the AEO release.\n\n"
        "The result is eleven EIA vintage releases, twenty forecast-outcome pairs, and all values "
        "expressed in real 2025 dollars throughout."
    ))

    # ── 6: EVALUATION DESIGN ──────────────────────────────────────────────
    slide = add_slide(prs, "Evaluation Design",
                      "Metrics and benchmark definition")
    # Left: metrics
    add_textbox(slide, Inches(0.72), Inches(1.42), Inches(5.85), Inches(0.45),
                "Forecast Error Definition", 18, NAVY, True)
    # Formula box (lightly tinted)
    _shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.72), Inches(1.92), Inches(5.85), Inches(0.72),
                TEAL_BG, TEAL,
                [("Forecast error  =  forecast − realized", 20, TEAL, True, PP_ALIGN.CENTER)],
                line_width_pt=1.5)
    add_bullets(slide, [
        "ME (Mean Error) — average bias; positive = optimistic",
        "MAE — mean absolute magnitude of error",
        "RMSE — root mean squared error",
        "MAPE — error as share of realized price",
    ], Inches(0.72), Inches(2.80), Inches(5.85), Inches(2.5), 17.5, NAVY, gap_before=6)
    # Vertical divider
    vdiv = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.88), Inches(1.42), Pt(1), Inches(4.5))
    vdiv.fill.solid()
    vdiv.fill.fore_color.rgb = LIGHT_GRAY
    vdiv.line.fill.background()
    # Right: random walk
    add_textbox(slide, Inches(7.15), Inches(1.42), Inches(5.5), Inches(0.45),
                "Random-Walk Benchmark", 18, NAVY, True)
    add_textbox(slide, Inches(7.15), Inches(1.97), Inches(5.5), Inches(0.82),
                "Forecast = most recently completed annual\nBrent average available at time of release.",
                18.5, NAVY, False, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(7.15), Inches(2.95), Inches(5.5), Inches(0.82),
                "Simple, transparent, and the standard commodity-price baseline.\n"
                "If EIA beats this, it demonstrates genuine informational value.",
                15.5, MID_GRAY, False, PP_ALIGN.LEFT)
    _shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7.15), Inches(4.25), Inches(5.5), Inches(0.60),
                BLUE_BG, BLUE,
                [("Simple · Transparent · Hard to beat", 16, BLUE, True, PP_ALIGN.CENTER)])
    set_notes(slide, (
        "Now let me explain our evaluation framework.\n\n"
        "The fundamental concept is simple: forecast error equals forecast minus realized. "
        "A positive error means the forecast was above the realized price — that is optimistic. "
        "A negative error means the forecast was below realized — that is pessimistic.\n\n"
        "From that base we compute four metrics. The Mean Error, or ME, is the average signed "
        "error. It tells us the direction and magnitude of bias. If ME is large and positive, "
        "forecasts are consistently too high on average. MAE and RMSE both measure average "
        "magnitude without sign cancellation — they tell us how far off forecasts typically are. "
        "And MAPE expresses error relative to realized price, giving a scale-free percentage.\n\n"
        "Now the benchmark. A random walk forecast means taking the most recently completed "
        "annual Brent price at the time of the AEO release and carrying it forward unchanged. "
        "No model. No analysis. Just the current level as the prediction.\n\n"
        "This is a meaningful benchmark for a specific reason. In commodity markets, current "
        "prices already reflect all publicly known information. The random walk hypothesis says "
        "that without private information or superior modeling, you cannot beat the current price "
        "as a forecast. If EIA beats this, the institutional modeling effort adds real value.\n\n"
        "And even if EIA beats this benchmark, we still need to check whether the forecasts are "
        "biased — because you can be informative and systematically wrong in one direction at "
        "the same time."
    ))

    # ── 7: EIA VINTAGES VS ACTUAL ─────────────────────────────────────────
    slide = add_slide(prs, "EIA Forecast Vintages vs Realized Brent",
                      "Institutional forecasts follow smooth paths; realized prices are far more volatile")
    add_picture_fit(slide,
                    PROJECT_ROOT / "outputs/eia_vintages_vs_actual.png",
                    Inches(0.60), Inches(1.35), Inches(12.15), Inches(4.92))
    add_textbox(slide, Inches(0.72), Inches(6.36), Inches(11.9), Inches(0.45),
                "Takeaway: EIA vintages are smooth, mean-reverting trend lines; "
                "realized Brent is highly volatile — motivating formal accuracy and bias evaluation.",
                15, TEAL, True, PP_ALIGN.CENTER)
    set_notes(slide, (
        "Before going to the numbers, I want to show you what these forecasts actually look like.\n\n"
        "This chart plots each EIA vintage as a forecast path — what EIA projected Brent would "
        "do from the release year forward. The thick black line is realized annual Brent.\n\n"
        "The contrast is immediately visible. Every EIA vintage follows a smooth, gradual "
        "trajectory — forecasts are mean-reverting and steady. Realized Brent, by contrast, is "
        "highly volatile. It collapsed in 2014 and 2015, spiked in 2022, and generally moved in "
        "ways that no smooth trend line could anticipate.\n\n"
        "Two things stand out. First, the forecast paths tend to sit above realized prices during "
        "extended low-price periods. That is the visual signature of the optimistic bias we will "
        "quantify on the next slide.\n\n"
        "Second, the smoothness of the forecasts reflects the nature of long-run equilibrium "
        "modeling. EIA is projecting where oil prices should settle over time, not where they "
        "will be next quarter. That is a deliberate modeling choice — but it means forecasts are "
        "structurally unable to capture the supply shocks and demand disruptions that actually "
        "drive short-to-medium-run oil prices.\n\n"
        "This chart motivates the formal analysis. The pattern is visible; now we need to "
        "measure it precisely."
    ))

    # ── 8: EIA BEATS RANDOM WALK ──────────────────────────────────────────
    slide = add_slide(prs, "EIA Beats the Random Walk",
                      "EIA forecasts contain real information relative to a carry-forward rule")
    eia3 = metric("EIA", 3, "RMSE")
    rw3 = metric("Random Walk", 3, "RMSE")
    eia5 = metric("EIA", 5, "RMSE")
    rw5 = metric("Random Walk", 5, "RMSE")
    eia_mae3 = metric("EIA", 3, "MAE")
    rw_mae3 = metric("Random Walk", 3, "MAE")
    eia_mae5 = metric("EIA", 5, "MAE")
    rw_mae5 = metric("Random Walk", 5, "MAE")

    # Two large comparison blocks
    for col, (horiz, ev, rv) in enumerate([("3-Year", eia3, rw3), ("5-Year", eia5, rw5)]):
        bx = Inches(0.78) + col * Inches(6.22)
        bw = Inches(5.75)
        half = Inches(2.42)
        # Horizon label
        add_textbox(slide, bx, Inches(1.44), bw, Inches(0.36),
                    f"{horiz} RMSE", 14, MID_GRAY, align=PP_ALIGN.CENTER)
        # EIA number (teal)
        add_textbox(slide, bx, Inches(1.78), half, Inches(0.30),
                    "EIA", 14, TEAL, True, PP_ALIGN.CENTER)
        add_textbox(slide, bx, Inches(2.06), half, Inches(0.88),
                    f"{ev:.2f}", 52, TEAL, True, PP_ALIGN.CENTER)
        # vs
        add_textbox(slide, bx + half, Inches(2.28), Inches(0.92), Inches(0.45),
                    "vs", 15, MID_GRAY, align=PP_ALIGN.CENTER)
        # Random Walk number (gray)
        add_textbox(slide, bx + half + Inches(0.92), Inches(1.78), half, Inches(0.30),
                    "Random Walk", 14, MID_GRAY, True, PP_ALIGN.CENTER)
        add_textbox(slide, bx + half + Inches(0.92), Inches(2.06), half, Inches(0.88),
                    f"{rv:.2f}", 52, MID_GRAY, True, PP_ALIGN.CENTER)
        # Improvement note
        diff = rv - ev
        pct = (diff / rv) * 100
        add_textbox(slide, bx, Inches(3.10), bw, Inches(0.36),
                    f"EIA is {diff:.2f} lower  ({pct:.0f}% reduction)",
                    15, TEAL, True, PP_ALIGN.CENTER)
        # Column divider after first block
        if col == 0:
            vd2 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(6.67), Inches(1.42), Pt(1), Inches(2.1))
            vd2.fill.solid()
            vd2.fill.fore_color.rgb = LIGHT_GRAY
            vd2.line.fill.background()

    add_divider(slide, Inches(0.72), Inches(3.62), Inches(11.9))
    eia_mae3 = metric("EIA", 3, "MAE")
    rw_mae3 = metric("Random Walk", 3, "MAE")
    eia_mae5 = metric("EIA", 5, "MAE")
    rw_mae5 = metric("Random Walk", 5, "MAE")
    add_bullets(slide, [
        (f"EIA also beats random walk on MAE: {eia_mae3:.2f} vs {rw_mae3:.2f} (3-year);  "
         f"{eia_mae5:.2f} vs {rw_mae5:.2f} (5-year)."),
        "Finding: EIA long-term Brent forecasts are informative — they outperform a simple carry-forward rule.",
    ], Inches(0.82), Inches(3.80), Inches(11.6), Inches(1.1), 18, NAVY, gap_before=5)
    set_notes(slide, (
        "Here is our first key finding.\n\n"
        f"At the three-year horizon, EIA's RMSE is {eia3:.2f} dollars per barrel, compared to "
        f"{rw3:.2f} for the random walk. EIA wins by about nineteen percent. "
        f"At the five-year horizon, EIA's RMSE is {eia5:.2f}, against {rw5:.2f} for the random "
        "walk — a twenty-eight percent reduction.\n\n"
        f"The same pattern holds for Mean Absolute Error. At three years, EIA's MAE is "
        f"{eia_mae3:.2f} against {rw_mae3:.2f} for the random walk. "
        f"At five years, {eia_mae5:.2f} versus {rw_mae5:.2f}.\n\n"
        "EIA beats the random walk cleanly, on both metrics, at both horizons. "
        "That is our first finding: EIA long-term Brent forecasts contain real informational value.\n\n"
        "This is not a trivial result. Random walk forecasts are famously hard to beat for "
        "commodity prices, precisely because the current price already incorporates all publicly "
        "available information. The fact that EIA consistently clears this bar at both horizons "
        "tells us that the institutional modeling effort is producing something genuinely useful.\n\n"
        "Practically, this means you cannot dismiss EIA forecasts as uninformative. There is a "
        "statistical case for using them as an input rather than ignoring them.\n\n"
        "But — and this is critical — beating the random walk on RMSE does not mean the forecasts "
        "are unbiased. You can outperform a naive benchmark in overall accuracy while still being "
        "consistently wrong in one direction. That is exactly what the next slide reveals."
    ))

    # ── 9: BIAS RESULT ────────────────────────────────────────────────────
    slide = add_slide(prs, "But EIA Is Optimistically Biased",
                      "Positive mean error indicates systematic overprediction of realized Brent")
    me3 = metric("EIA", 3, "ME")
    me5 = metric("EIA", 5, "ME")

    # Two large ME callouts (left 2/3 of slide)
    for col, (horiz, me) in enumerate([("3-Year Mean Error", me3), ("5-Year Mean Error", me5)]):
        cx = Inches(0.72) + col * Inches(4.2)
        add_textbox(slide, cx, Inches(1.46), Inches(3.85), Inches(0.32),
                    horiz, 14, MID_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx, Inches(1.76), Inches(3.85), Inches(1.1),
                    f"+{me:.2f}", 62, ORANGE, True, PP_ALIGN.CENTER)
        add_textbox(slide, cx, Inches(2.84), Inches(3.85), Inches(0.36),
                    "USD/bbl average overprediction", 14, MID_GRAY, align=PP_ALIGN.CENTER)

    add_divider(slide, Inches(0.72), Inches(3.36), Inches(7.65))
    add_bullets(slide, [
        "Positive ME means EIA forecasts are above realized Brent on average — optimistically biased.",
        "Bias compounds in DCF: inflated price → inflated revenue → inflated NPV and IRR.",
        "Implication: EIA is informative, but should not be used mechanically without bias adjustment.",
    ], Inches(0.82), Inches(3.55), Inches(7.55), Inches(2.15), 18, NAVY, gap_before=7)

    # Risk callout box (right side)
    _shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.82), Inches(1.42), Inches(3.85), Inches(4.35),
                ORANGE_BG, ORANGE,
                [
                    ("Capital Budgeting Risk", 16, ORANGE, True, PP_ALIGN.CENTER),
                    ("\nHigher forecast prices inflate projected revenue, NPV, and IRR — "
                     "creating a systematic bias toward over-approval of upstream projects.",
                     14, NAVY, False, PP_ALIGN.CENTER, 0, 1.2),
                    ("\nRisk: Overinvestment", 15, RED, True, PP_ALIGN.CENTER, 8),
                ],
                line_width_pt=1.5, margin_h=Inches(0.18), margin_v=Inches(0.18))
    set_notes(slide, (
        "This is, in my view, the most important finding in the study.\n\n"
        f"At the three-year horizon, EIA's mean error is plus {me3:.2f} dollars per barrel. "
        f"At the five-year horizon, it is plus {me5:.2f}.\n\n"
        "Let me make sure that is clear. These are signed averages, not magnitudes. "
        "Positive means that on average, EIA forecasts sat roughly twenty to twenty-seven "
        "dollars above what actually happened. That is a substantial, persistent, directional "
        "bias toward optimism.\n\n"
        "To put twenty dollars per barrel in context: on a project producing one million barrels "
        "per year, that is twenty million dollars of annual revenue that exists only in the model, "
        "not in the world. Over a ten-year project life, you are looking at hundreds of millions "
        "of dollars of phantom revenue in the DCF. That is enough to push borderline projects "
        "over the hurdle rate.\n\n"
        "The risk box on the right captures the capital budgeting consequence. Higher forecast "
        "prices generate inflated revenue projections. Inflated revenues feed into the DCF. "
        "NPV and IRR both look better. Projects that should fail the hurdle rate instead pass it. "
        "That is a systematic bias toward over-approval of upstream investment.\n\n"
        "To be precise about what we have found: EIA forecasts are both informative and biased. "
        "They beat the random walk — that is the informative part. But they are consistently "
        f"above realized prices by {me3:.0f} to {me5:.0f} dollars — that is the bias. "
        "Both findings are real, and the combination tells you exactly how to use these forecasts: "
        "with an adjustment, not as-is."
    ))

    # ── 10: ROLLING MODEL EXTENSION ───────────────────────────────────────
    slide = add_slide(prs, "Rolling Model Extension",
                      "Simple econometric extensions do not improve on EIA")
    ar3 = metric("Rolling AR", 3, "RMSE", rolling)
    ar5 = metric("Rolling AR", 5, "RMSE", rolling)
    arf3 = metric("Rolling AR + Futures", 3, "RMSE", rolling)
    arf5 = metric("Rolling AR + Futures", 5, "RMSE", rolling)
    rolling_data = [
        ["Model", "3-Year RMSE", "5-Year RMSE"],
        ["EIA", f"{eia3:.2f}", f"{eia5:.2f}"],
        ["Rolling AR", f"{ar3:.2f}", f"{ar5:.2f}"],
        ["Rolling AR + Futures", f"{arf3:.2f}", f"{arf5:.2f}"],
    ]
    add_table(slide, rolling_data,
              Inches(0.78), Inches(1.44), Inches(7.0), Inches(2.85),
              [Inches(3.05), Inches(1.88), Inches(1.88)], 17)

    # Right: model detail bullets
    add_textbox(slide, Inches(8.18), Inches(1.44), Inches(4.5), Inches(0.45),
                "Model Details", 16, NAVY, True)
    add_bullets(slide, [
        "Rolling AR: autoregression on lagged annual Brent",
        "AR + Futures: adds front-month BZ=F futures and basis spread",
        "Estimated on expanding window; predicted out-of-sample",
    ], Inches(8.18), Inches(1.95), Inches(4.5), Inches(1.92), 16, NAVY, gap_before=7)

    add_divider(slide, Inches(0.78), Inches(4.52), Inches(11.9))
    add_textbox(slide, Inches(0.82), Inches(4.70), Inches(11.6), Inches(0.55),
                "Finding: Neither rolling extension beats EIA. EIA remains the best available forecast in this comparison.",
                19, TEAL, True, PP_ALIGN.CENTER)
    set_notes(slide, (
        "A natural question is whether we can build a better forecast. We tested two extensions.\n\n"
        "The first is a rolling autoregressive model using lagged annual Brent prices estimated "
        "on an expanding window, predicting the next three- or five-year average out of sample. "
        "The second adds front-month BZ=F futures prices and the basis spread as additional "
        "regressors.\n\n"
        "The results are in the table. "
        f"At three years, EIA's RMSE of {eia3:.2f} beats the rolling AR at {ar3:.2f} "
        f"and the AR plus futures at {arf3:.2f}. "
        f"At five years, EIA's {eia5:.2f} beats the rolling AR at {ar5:.2f} "
        f"and the AR plus futures at {arf5:.2f}.\n\n"
        "Neither extension beats EIA. EIA remains the most accurate model in this comparison.\n\n"
        "A few caveats. Our rolling models are intentionally simple. The futures variable is a "
        "front-month contract, which is a rough proxy for long-horizon price expectations — it "
        "is not a forward curve. More sophisticated models using longer-dated futures or richer "
        "information sets might perform differently. And our out-of-sample window is limited by "
        "the sample size.\n\n"
        "For our purposes: EIA is the best available baseline in this comparison, and the simple "
        "extensions we tested do not change that conclusion."
    ))

    # ── 11: INTERPRETATION ────────────────────────────────────────────────
    slide = add_slide(prs, "Interpretation", "Useful, but not unbiased")
    # Vertical divider
    vdiv2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.05), Inches(1.42), Pt(1), Inches(4.2))
    vdiv2.fill.solid()
    vdiv2.fill.fore_color.rgb = LIGHT_GRAY
    vdiv2.line.fill.background()
    # Left column
    add_textbox(slide, Inches(0.72), Inches(1.42), Inches(6.0), Inches(0.48),
                "Informative relative to naive benchmark", 20, TEAL, True)
    add_bullets(slide, [
        "EIA beats random walk on RMSE and MAE at both horizons.",
        "Institutional forecasts add real informational content.",
        "Dismissing EIA forecasts entirely would be premature.",
    ], Inches(0.72), Inches(2.05), Inches(6.0), Inches(2.2), 18, NAVY, gap_before=7)
    # Right column
    add_textbox(slide, Inches(7.38), Inches(1.42), Inches(5.3), Inches(0.48),
                "Biased enough to matter for investment", 20, ORANGE, True)
    add_bullets(slide, [
        f"Consistent positive ME of ${me3:.0f}–${me5:.0f}/bbl at 3- and 5-year horizons.",
        "Bias flows directly into revenue and valuation models.",
        "Mechanical use of EIA price decks risks overinvestment.",
    ], Inches(7.38), Inches(2.05), Inches(5.3), Inches(2.2), 18, NAVY, gap_before=7)
    add_divider(slide, Inches(0.72), Inches(4.78), Inches(11.9))
    add_textbox(slide, Inches(0.72), Inches(4.96), Inches(11.9), Inches(0.55),
                "Synthesis: EIA is informative but biased — use it as a starting point, not a final answer.",
                20, NAVY, True, PP_ALIGN.CENTER)
    set_notes(slide, (
        "Let me step back and synthesize the findings before moving to implications.\n\n"
        "The results split into two clean messages.\n\n"
        "On the left: EIA is informative relative to a naive benchmark. It beats the random walk "
        "on RMSE and MAE at both horizons. That is real statistical evidence that institutional "
        "forecasting adds value. It means there is a genuine empirical case for using these "
        "forecasts as inputs rather than discarding them.\n\n"
        "On the right: the forecasts are biased enough to matter for investment decisions. "
        f"A mean error of {me3:.0f} to {me5:.0f} dollars per barrel is not a rounding error. "
        "Over the life of a major upstream project, that level of systematic optimism can "
        "materially distort NPV calculations and lead to approval of projects that would not "
        "pass a realistic price test.\n\n"
        "These two results coexist. The forecast is informative and biased simultaneously.\n\n"
        "The management response has to reflect both parts. Do not discard EIA forecasts — that "
        "discards real signal. Do not use them mechanically — that embeds the bias directly into "
        "capital budgeting.\n\n"
        "The correct approach is to use EIA as a starting point, then apply a systematic "
        "correction. That is what the next slide recommends."
    ))

    # ── 12: MANAGERIAL IMPLICATIONS ───────────────────────────────────────
    slide = add_slide(prs, "Managerial Implications",
                      "How upstream firms should use institutional price decks")
    implications = [
        ("01", TEAL,
         "Use EIA as a baseline, not a final answer",
         "EIA beats the random walk — discarding it entirely would discard genuine information."),
        ("02", TEAL,
         "Apply a bias adjustment or conservative haircut",
         f"Historical ME: –${me3:.0f}/bbl at 3-year; –${me5:.0f}/bbl at 5-year. Run sensitivities around this."),
        ("03", ORANGE,
         "Run downside price scenarios before final investment decisions",
         "Stress-test projects against plausible realized price paths, not smooth forecast paths."),
        ("04", ORANGE,
         "Treat long-horizon price uncertainty as a capital allocation risk",
         "Forecast uncertainty should appear in project risk registers and inform hurdle-rate calibration."),
    ]
    for i, (num, col, heading, detail) in enumerate(implications):
        y = Inches(1.48) + i * Inches(1.22)
        # Badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), y + Inches(0.06),
            Inches(0.52), Inches(0.52))
        badge.fill.solid()
        badge.fill.fore_color.rgb = col
        badge.line.fill.background()
        btf = badge.text_frame
        _tf_setup(btf, margin_h=Inches(0.0), margin_v=Inches(0.0))
        _para_run(btf, num, 13, WHITE, True, PP_ALIGN.CENTER, first=True)
        # Heading
        add_textbox(slide, Inches(1.44), y, Inches(11.2), Inches(0.44),
                    heading, 19, NAVY, True, PP_ALIGN.LEFT)
        # Detail
        add_textbox(slide, Inches(1.44), y + Inches(0.44), Inches(11.2), Inches(0.38),
                    detail, 14.5, MID_GRAY, False, PP_ALIGN.LEFT)
        if i < 3:
            add_divider(slide, Inches(0.72), y + Inches(0.92), Inches(11.9))
    set_notes(slide, (
        "So what should upstream firms actually do?\n\n"
        "Four practical recommendations.\n\n"
        "First: use EIA as a baseline. It carries real informational content — it beats the "
        "random walk. Use it as the starting point for your price path. Do not build price "
        "decks from scratch, and do not ignore the institutional forecast.\n\n"
        "Second: apply a bias adjustment. Our results give you direct empirical guidance. "
        f"At the three-year horizon, the historical mean error is around minus "
        f"{me3:.0f} dollars per barrel. At five years, around minus {me5:.0f} dollars. "
        "That is the average correction needed to align EIA forecasts with realized prices. "
        "Apply it as a haircut to the baseline path, and run sensitivity analysis around it.\n\n"
        "Third: run downside scenarios. Bias adjustment corrects the average, but there is still "
        "substantial uncertainty around that average. Before sanctioning any major project, test "
        "it explicitly against a price path that reflects historical realized volatility — not "
        "just against a smooth, corrected trend line.\n\n"
        "Fourth: treat long-horizon price uncertainty as a capital allocation risk item. This "
        "uncertainty should appear explicitly in the project risk register, not be absorbed "
        "silently into a single deterministic price assumption. And it should influence how the "
        "organization calibrates its hurdle rate.\n\n"
        "The overall recommendation: adjust institutional forecasts rather than either "
        "discarding or mechanically applying them."
    ))

    # ── 13: LIMITATIONS AND NEXT STEPS ───────────────────────────────────
    slide = add_slide(prs, "Limitations and Next Steps",
                      "Scope boundaries and natural extensions")
    vdiv3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.05), Inches(1.42), Pt(1), Inches(5.15))
    vdiv3.fill.solid()
    vdiv3.fill.fore_color.rgb = LIGHT_GRAY
    vdiv3.line.fill.background()
    # Limitations
    add_textbox(slide, Inches(0.72), Inches(1.42), Inches(6.0), Inches(0.45),
                "Limitations", 20, NAVY, True)
    add_bullets(slide, [
        "Main analysis covers EIA; IEA and other institutions not compared.",
        "World Bank uses Crude oil (avg), not Brent spot — different contract.",
        "Futures extension uses front-month BZ=F, not a long-dated forward curve.",
        "Long-horizon realized sample is naturally limited (~10–11 outcomes).",
    ], Inches(0.72), Inches(2.02), Inches(6.0), Inches(3.55), 17, NAVY, gap_before=7)
    # Next steps
    add_textbox(slide, Inches(7.38), Inches(1.42), Inches(5.6), Inches(0.45),
                "Next Steps", 20, NAVY, True)
    add_bullets(slide, [
        "Add IEA or World Bank as multi-institution robustness check.",
        "Test bias-corrected EIA: apply historical ME adjustment to future vintages.",
        "Use long-dated Brent futures curve if reliable historical data are available.",
        "Extend panel as new EIA AEO vintages are published annually.",
    ], Inches(7.38), Inches(2.02), Inches(5.6), Inches(3.55), 17, NAVY, gap_before=7)
    set_notes(slide, (
        "Let me be transparent about the scope boundaries of this study.\n\n"
        "Our main analysis focuses on EIA. We have not compared it against the IEA or other "
        "institutional forecasters, so we cannot assess relative performance across institutions — "
        "only how EIA performs against a random walk.\n\n"
        "The World Bank crude oil price series tracks a composite crude average rather than Brent "
        "specifically, which is why we keep it as a supplementary check rather than a core "
        "comparison point.\n\n"
        "Our futures extension uses a front-month contract — a rough proxy. A proper "
        "futures-based model would require longer-dated contracts that were consistently "
        "available over the sample period.\n\n"
        "And at five-year horizons, our realized sample naturally contains only about ten to "
        "eleven outcome observations. That limits statistical precision and means we should "
        "interpret the results with appropriate caution.\n\n"
        "The most compelling next extension is bias-corrected EIA forecasts: applying the "
        "historical mean error as a systematic correction and testing whether adjusted forecasts "
        "outperform raw EIA out of sample. That follows directly from our main finding."
    ))

    # ── 14: CONCLUSION ────────────────────────────────────────────────────
    slide = add_slide(prs, "Conclusion")
    # Main conclusion sentence
    _shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.72), Inches(1.45), Inches(11.9), Inches(1.62),
                TEAL_BG, TEAL,
                [("EIA long-term Brent forecasts are informative, but optimistically biased; "
                  "upstream firms should adjust them before using them in capital budgeting.",
                  21, NAVY, True, PP_ALIGN.CENTER, 0, 1.3)],
                line_width_pt=2.0, margin_h=Inches(0.28), margin_v=Inches(0.18))
    # Three takeaway chips
    takeaways = [
        ("Informative", "EIA beats random walk — use it as a baseline", TEAL),
        ("Biased", f"ME of +${me3:.0f}–${me5:.0f}/bbl — apply a haircut", ORANGE),
        ("Actionable", "Adjust before DCF — run downside scenarios", BLUE),
    ]
    chip_w = Inches(3.75)
    for i, (head, body, col) in enumerate(takeaways):
        cx = Inches(0.72) + i * Inches(4.05)
        _shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    cx, Inches(3.38), chip_w, Inches(1.42),
                    WHITE, col,
                    [
                        (head, 19, col, True, PP_ALIGN.CENTER),
                        (body, 14, NAVY, False, PP_ALIGN.CENTER, 5, 1.1),
                    ],
                    line_width_pt=2.0, margin_h=Inches(0.14), margin_v=Inches(0.12))
    add_textbox(slide, Inches(0.72), Inches(5.02), Inches(11.9), Inches(0.52),
                "Use institutional forecasts as starting points, not unadjusted truth.",
                19, TEAL, True, PP_ALIGN.CENTER)
    set_notes(slide, (
        "Let me close with the bottom line.\n\n"
        "EIA long-term Brent forecasts are informative, but optimistically biased. "
        "Upstream firms should adjust them before using them in capital budgeting.\n\n"
        "Three takeaways. EIA beats the random walk — use it as a baseline, not a blank sheet. "
        f"But with mean errors of {me3:.0f} to {me5:.0f} dollars per barrel, the bias is too "
        "large to ignore — apply a haircut before it enters the DCF. And treat price uncertainty "
        "as a risk management input, not a background assumption.\n\n"
        "The broader point is this: no long-horizon commodity price forecast is reliable enough "
        "to be used as truth. The best practice is to use institutional forecasts as informed "
        "starting points — acknowledge their limitations, correct for their biases, and "
        "stress-test against downside scenarios before capital is committed.\n\n"
        "Thank you. I am happy to take questions."
    ))

    # ── APPENDIX A ────────────────────────────────────────────────────────
    slide = add_slide(prs, "Appendix A: Full Benchmark Metrics")
    app_a = [["Model", "Horizon", "N", "ME", "MAE", "RMSE", "MAPE"]]
    for _, row in benchmark.iterrows():
        app_a.append([
            row["model"], int(row["horizon_years"]), int(row["n_forecasts"]),
            f"{row['ME']:.2f}", f"{row['MAE']:.2f}",
            f"{row['RMSE']:.2f}", f"{row['MAPE']:.3f}",
        ])
    add_table(slide, app_a, Inches(0.66), Inches(1.55), Inches(12.0), Inches(3.4),
              [Inches(2.2), Inches(1.25), Inches(0.75),
               Inches(1.45), Inches(1.45), Inches(1.45), Inches(1.25)], 13.5)
    set_notes(slide, (
        "Appendix A. This table shows the full benchmark metric results for both EIA and the "
        "random walk at three-year and five-year horizons. ME, MAE, RMSE, and MAPE are all "
        "reported. Refer to this if the audience asks for numbers beyond what was shown on the "
        "main result slides."
    ))

    # ── APPENDIX B ────────────────────────────────────────────────────────
    slide = add_slide(prs, "Appendix B: Rolling Model Metrics")
    app_b = [["Model", "Horizon", "N", "ME", "MAE", "RMSE", "MAPE"]]
    for _, row in rolling.iterrows():
        app_b.append([
            row["model"], int(row["horizon_years"]), int(row["n_forecasts"]),
            f"{row['ME']:.2f}", f"{row['MAE']:.2f}",
            f"{row['RMSE']:.2f}", f"{row['MAPE']:.3f}",
        ])
    add_table(slide, app_b, Inches(0.66), Inches(1.55), Inches(12.0), Inches(3.4),
              [Inches(2.55), Inches(1.15), Inches(0.65),
               Inches(1.35), Inches(1.35), Inches(1.35), Inches(1.15)], 13.5)
    set_notes(slide, (
        "Appendix B. Full metrics for the rolling AR and rolling AR plus futures models. "
        "Refer to this if the audience asks for the complete rolling model performance numbers "
        "including ME, MAE, and MAPE."
    ))

    # ── APPENDIX C ────────────────────────────────────────────────────────
    slide = add_slide(prs, "Appendix C: Data Source Links")
    add_bullets(slide, [
        "EIA Annual Energy Outlook Table 12 — eia.gov/aeo",
        "FRED DCOILBRENTEU: Brent Spot Price — fred.stlouisfed.org",
        "FRED CPIAUCSL: Consumer Price Index — fred.stlouisfed.org",
        "Yahoo Finance BZ=F: Brent crude front-month futures proxy",
        "World Bank Commodity Markets price archive — worldbank.org",
    ], Inches(1.0), Inches(1.70), Inches(10.8), Inches(3.6), 18, NAVY, gap_before=10)
    set_notes(slide, (
        "Appendix C. Data source links. Refer to this slide if the audience asks where the "
        "data came from or how to access the raw series used in the analysis."
    ))

    # ── APPENDIX D ────────────────────────────────────────────────────────
    slide = add_slide(prs, "Appendix D: Reproducibility")
    add_bullets(slide, [
        "Structured Python package: src/oil_forecast_project",
        "Pipeline writes CSVs, charts, Excel workbook, and summary automatically",
        "Tests cover metrics, CPI conversion, Brent aggregation, panel construction, and benchmark summary",
        "Validation: pytest passes with 6 tests",
    ], Inches(1.0), Inches(1.70), Inches(10.8), Inches(3.6), 18, NAVY, gap_before=10)
    set_notes(slide, (
        "Appendix D. Reproducibility details. The full analysis is implemented as a structured "
        "Python package under src/oil_forecast_project. All six tests pass. Refer to this slide "
        "if the audience asks about how to replicate the results."
    ))

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(f"Wrote {path}")
